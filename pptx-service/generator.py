"""PPTX generator — code-based design, zero template dependency for MVP.

Design theme: "Corporate" — Navy + Gold, Calibri, geometric accents.
Templates are code constants; forward-compatible with .pptx template files later.
"""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

from schemas import (
    Bullet,
    ChartSlide,
    ClosingSlide,
    ComparisonSlide,
    ContentSlide,
    CoverSlide,
    GenerateRequest,
    SectionDividerSlide,
    SlideMeta,
)

# ═══════════════════════════════════════════════
#  Design Tokens — "Corporate"
# ═══════════════════════════════════════════════

PRIMARY = RGBColor(0x1A, 0x36, 0x5D)       # Navy
SECONDARY = RGBColor(0x2B, 0x6C, 0xB0)      # Blue
ACCENT = RGBColor(0xED, 0x89, 0x36)         # Gold
BACKGROUND = RGBColor(0xF7, 0xFA, 0xFC)      # Off-white
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x2D, 0x37, 0x48)            # Dark gray
MUTED = RGBColor(0xA0, 0xAE, 0xC0)           # Light gray
ACCENT_LIGHT = RGBColor(0xCC, 0xD5, 0xE0)    # Faded navy (for subtitle on dark bg)

HEADING_FONT = "Calibri Light"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN * 2  # ~11.7 inches

# ═══════════════════════════════════════════════
#  Shape Helpers
# ═══════════════════════════════════════════════


def _rect(slide, left, top, width, height, fill_color=None, name=""):
    """Add a rectangle with optional solid fill."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    if name:
        shape.name = name
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height, name, text="",
             font_size=None, color=None, bold=False, align=PP_ALIGN.LEFT,
             font_name=None, anchor=MSO_ANCHOR.TOP):
    """Add a text box with single-paragraph text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size or Pt(18)
    p.font.color.rgb = color or TEXT
    p.font.bold = bold
    p.font.name = font_name or BODY_FONT
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def _set_para(para, text, size, color, bold=False, font_name=None, align=PP_ALIGN.LEFT):
    """Configure a paragraph with text and formatting."""
    para.text = text
    para.font.size = size
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.name = font_name or BODY_FONT
    para.alignment = align
    para.space_before = Pt(0)
    para.space_after = Pt(0)


def _bullet_frame(shape, bullets: list[Bullet]):
    """Fill a text frame with indented bullets."""
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b.text
        p.level = b.level
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT
        p.font.name = BODY_FONT
        p.space_before = Pt(4)
        p.space_after = Pt(4)


# ═══════════════════════════════════════════════
#  Slide Renderers
# ═══════════════════════════════════════════════


def _render_cover(slide, data: CoverSlide, meta: SlideMeta):
    # Navy block (left 55%)
    _rect(slide, Inches(0), Inches(0), Inches(7.3), SLIDE_H, PRIMARY)
    # Gold accent stripe
    _rect(slide, Inches(7.3), Inches(0), Inches(0.08), SLIDE_H, ACCENT)
    # Corner square decoration
    _rect(slide, Inches(6.8), Inches(0.5), Inches(0.35), Inches(0.35), ACCENT)

    # Title
    _textbox(slide, Inches(1), Inches(2), Inches(5.8), Inches(1.5),
             "title", data.title, Pt(46), WHITE, True, PP_ALIGN.LEFT, HEADING_FONT)
    # Subtitle
    if data.subtitle:
        _textbox(slide, Inches(1), Inches(3.7), Inches(5.8), Inches(0.7),
                 "subtitle", data.subtitle, Pt(22), ACCENT_LIGHT, False, PP_ALIGN.LEFT)
    # Date
    if data.date:
        _textbox(slide, Inches(1), Inches(4.7), Inches(3), Inches(0.5),
                 "date", data.date, Pt(14), MUTED, False)
    # Presenter
    if data.presenter:
        _textbox(slide, Inches(1), Inches(5.2), Inches(3), Inches(0.5),
                 "presenter", data.presenter, Pt(14), MUTED, False)

    # Right-side decorative line
    _rect(slide, Inches(8.5), Inches(1.5), Inches(3.5), Inches(0.015), ACCENT)


def _render_content(slide, data: ContentSlide):
    # Top accent bar
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), PRIMARY)
    # Title
    _textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "title", data.title, Pt(32), PRIMARY, True, PP_ALIGN.LEFT, HEADING_FONT)
    # Gold underline
    _rect(slide, MARGIN, Inches(1.25), Inches(2), Inches(0.035), ACCENT)
    # Body bullets
    body = _textbox(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(4.8),
                    "body", font_size=Pt(18), color=TEXT)
    _bullet_frame(body, data.bullets)
    # Footer line
    _rect(slide, MARGIN, Inches(6.9), CONTENT_W, Inches(0.008), MUTED)


def _render_section_divider(slide, data: SectionDividerSlide):
    # Full navy background
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
    # Gold left border
    _rect(slide, Inches(0), Inches(0), Inches(0.1), SLIDE_H, ACCENT)
    # Section number
    _textbox(slide, Inches(1.5), Inches(2), Inches(3), Inches(1.2),
             "number", data.section_number, Pt(72), ACCENT, True, PP_ALIGN.LEFT, HEADING_FONT)
    # Section title
    _textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2),
             "title", data.title, Pt(42), WHITE, True, PP_ALIGN.LEFT, HEADING_FONT)
    # Subtitle
    if data.subtitle:
        _textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.7),
                 "subtitle", data.subtitle, Pt(20), ACCENT_LIGHT, False)
    # Corner decoration
    _rect(slide, Inches(12.2), Inches(0.3), Inches(0.5), Inches(0.5), ACCENT)


def _render_comparison(slide, data: ComparisonSlide):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), PRIMARY)
    _textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "title", data.title, Pt(32), PRIMARY, True, PP_ALIGN.LEFT, HEADING_FONT)
    _rect(slide, MARGIN, Inches(1.25), Inches(2), Inches(0.035), ACCENT)

    left_x = MARGIN
    right_x = Inches(7.0)
    col_w = Inches(5.6)

    # ── Left column ──
    _rect(slide, left_x, Inches(1.7), col_w, Inches(0.05), SECONDARY)
    _textbox(slide, left_x, Inches(1.9), col_w, Inches(0.6),
             "left_h", data.left.heading, Pt(22), SECONDARY, True, font_name=HEADING_FONT)
    lb = _textbox(slide, left_x, Inches(2.6), col_w, Inches(4), "left_body")
    tf = lb.text_frame
    tf.clear()
    for i, pt in enumerate(data.left.points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, f"• {pt}", Pt(16), TEXT)

    # ── Right column ──
    _rect(slide, right_x, Inches(1.7), col_w, Inches(0.05), SECONDARY)
    _textbox(slide, right_x, Inches(1.9), col_w, Inches(0.6),
             "right_h", data.right.heading, Pt(22), SECONDARY, True, font_name=HEADING_FONT)
    rb = _textbox(slide, right_x, Inches(2.6), col_w, Inches(4), "right_body")
    tf = rb.text_frame
    tf.clear()
    for i, pt in enumerate(data.right.points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, f"• {pt}", Pt(16), TEXT)


def _render_chart(slide, data: ChartSlide):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), PRIMARY)
    _textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
             "title", data.title, Pt(32), PRIMARY, True, PP_ALIGN.LEFT, HEADING_FONT)
    _rect(slide, MARGIN, Inches(1.25), Inches(2), Inches(0.035), ACCENT)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = data.categories
    for s in data.series:
        chart_data.add_series(s.name, s.values)

    chart_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }[data.chart_type]

    chart_frame = slide.shapes.add_chart(
        chart_type, Inches(0.8), Inches(1.7), Inches(8.5), Inches(4.8), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # Style the chart series colors
    series_colors = [PRIMARY, SECONDARY, ACCENT, MUTED]
    for i, s in enumerate(chart.series):
        if i < len(series_colors):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = series_colors[i]

    # Insight callout (dark box on right)
    if data.insight:
        _rect(slide, Inches(9.8), Inches(1.7), Inches(2.8), Inches(1.8), PRIMARY)
        _textbox(slide, Inches(10.0), Inches(1.85), Inches(2.4), Inches(1.5),
                 "insight", data.insight, Pt(14), WHITE, False)


def _render_closing(slide, data: ClosingSlide, meta: SlideMeta):
    # Top navy block
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.5), PRIMARY)
    # Separator
    _rect(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.05), ACCENT)
    # Title
    _textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1),
             "title", data.title, Pt(44), WHITE, True, PP_ALIGN.LEFT, HEADING_FONT)
    # Subtitle
    if data.subtitle:
        _textbox(slide, Inches(1.5), Inches(2), Inches(10), Inches(0.7),
                 "subtitle", data.subtitle, Pt(22), ACCENT_LIGHT, False)
    # Contact
    if data.contact:
        _textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.7),
                 "contact", data.contact, Pt(16), MUTED, False)
    # Decorative squares
    _rect(slide, Inches(12.2), Inches(6.5), Inches(0.7), Inches(0.7), ACCENT)
    _rect(slide, Inches(12.4), Inches(6.7), Inches(0.7), Inches(0.7), SECONDARY)


RENDERERS = {
    "cover": _render_cover,
    "content": _render_content,
    "section_divider": _render_section_divider,
    "comparison": _render_comparison,
    "chart": _render_chart,
    "closing": _render_closing,
}

# slide types that need meta context
NEEDS_META = {"cover", "closing"}


# ═══════════════════════════════════════════════
#  Main Generator
# ═══════════════════════════════════════════════

def generate(request: GenerateRequest) -> bytes:
    """Generate a .pptx file from a GenerateRequest. Returns bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Set background on the default blank layout
    blank_layout = prs.slide_layouts[6]  # blank layout (index 6 in default template)

    # Remove any default slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))  # noqa: SLF001
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])  # noqa: SLF001

    section_num = 0
    for slide_data in request.slides:
        # Auto-fill section_number for section dividers
        if slide_data.type == "section_divider" and not slide_data.section_number:
            section_num += 1
            slide_data.section_number = str(section_num).zfill(2)

        slide = prs.slides.add_slide(blank_layout)

        # Set slide background to off-white
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = BACKGROUND

        renderer = RENDERERS.get(slide_data.type)
        if renderer:
            if slide_data.type in NEEDS_META:
                renderer(slide, slide_data, request.meta)
            else:
                renderer(slide, slide_data)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
